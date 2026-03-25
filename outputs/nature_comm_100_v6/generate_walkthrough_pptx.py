"""Generate CrossDisc_Bench pipeline walkthrough PPT using v6 data."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import json, os

BASE = os.path.dirname(os.path.abspath(__file__))

# ── Load v6 data ──
with open(os.path.join(BASE, "extraction_results.json"), encoding="utf-8") as f:
    extraction = json.load(f)
with open(os.path.join(BASE, "benchmark_dataset.json"), encoding="utf-8") as f:
    benchmark = json.load(f)
with open(os.path.join(BASE, "p1p5_comparison_summary.json"), encoding="utf-8") as f:
    p1p5 = json.load(f)
with open(os.path.join(BASE, "p5_kg_eval_results.json"), encoding="utf-8") as f:
    kg_eval = json.load(f)

paper0 = extraction[0]
gt0 = benchmark[0]
kg0 = kg_eval[0] if kg_eval else None

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color constants ──
BLUE = RGBColor(0x1A, 0x56, 0xDB)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
CODE_BG = RGBColor(0xF2, 0xF2, 0xF2)
GREEN = RGBColor(0x0D, 0x92, 0x76)
ORANGE = RGBColor(0xE8, 0x6C, 0x00)

# ── Helper functions ──

def add_slide(title_text, layout_idx=1):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.font.name = "\u5fae\u8f6f\u96c5\u9ed1"
    line = slide.shapes.add_shape(1, Inches(0.8), Inches(1.15), Inches(11.5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()
    return slide

def add_text_box(slide, left, top, width, height, text, font_size=16,
                 bold=False, color=DARK, font_name="\u5fae\u8f6f\u96c5\u9ed1", align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf

def add_bullet_list(slide, left, top, width, height, items, font_size=15, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"\u2022 {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK
        p.font.name = "\u5fae\u8f6f\u96c5\u9ed1"
        p.space_after = spacing
    return tf

def add_code_box(slide, left, top, width, height, text, font_size=11):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(8)
    tf.margin_right = Pt(12)
    tf.margin_bottom = Pt(8)
    for i, line in enumerate(text.strip().split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.name = "Consolas"
        p.font.color.rgb = DARK
    return tf

def add_card(slide, left, top, width, height, title, content, title_color=BLUE, bg_color=None):
    if bg_color:
        bg = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_color
        bg.line.fill.background()
    # Title bar
    bar = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(0.45))
    bar.fill.solid()
    bar.fill.fore_color.rgb = title_color
    bar.line.fill.background()
    tf_t = bar.text_frame
    tf_t.margin_left = Pt(10)
    tf_t.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf_t.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "\u5fae\u8f6f\u96c5\u9ed1"
    # Content
    if content:
        txBox = slide.shapes.add_textbox(
            Inches(left + 0.15), Inches(top + 0.55),
            Inches(width - 0.3), Inches(height - 0.65))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(12)
        p.font.color.rgb = DARK
        p.font.name = "\u5fae\u8f6f\u96c5\u9ed1"
        return tf
    return None

def add_table_slide(slide, left, top, width, height, headers, rows):
    table_shape = slide.shapes.add_table(
        len(rows) + 1, len(headers),
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    table = table_shape.table
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "\u5fae\u8f6f\u96c5\u9ed1"
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = "\u5fae\u8f6f\u96c5\u9ed1"
                p.font.color.rgb = DARK
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG
    return table

# Compute stats
n_papers = len(extraction)
n_gt = len(benchmark)
all_concepts_counts = [sum(len(v) for v in p.get("concepts", {}).values()) for p in extraction]
all_relations_counts = [len(p.get("relations", [])) for p in extraction]
avg_c = sum(all_concepts_counts) / max(len(all_concepts_counts), 1)
avg_r = sum(all_relations_counts) / max(len(all_relations_counts), 1)
all_gt_terms = [p.get("gt_stats", {}).get("n_terms", 0) for p in benchmark]
all_gt_rels = [p.get("gt_stats", {}).get("n_relations", 0) for p in benchmark]

title = paper0.get("title", "")
primary = paper0.get("primary", "")
secondary = paper0.get("secondary", "")

# ══════════════════════════════════════════════════════════════════════
# Slide 1: Title
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(0x0A, 0x1A, 0x3A)
bg.line.fill.background()

add_text_box(slide, 1.5, 1.8, 10, 1.2,
             "CrossDisc-Bench Pipeline Walkthrough",
             font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 3.0, 10, 0.8,
             "\u8de8\u5b66\u79d1\u5047\u8bbe\u751f\u6210\u57fa\u51c6\u6d4b\u8bd5 \u2014 \u5168\u6d41\u7a0b\u5b9e\u4f8b\u5c55\u793a (v6)",
             font_size=22, color=RGBColor(0xBB, 0xCC, 0xEE), align=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 4.2, 10, 0.6,
             f"Nature Communications | {n_papers} Papers | {n_gt} GT Entries",
             font_size=18, color=GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════
# Slide 2: Pipeline Overview
# ══════════════════════════════════════════════════════════════════════
slide = add_slide("1. Pipeline Overview (\u6d41\u7a0b\u603b\u89c8)")

stages = [
    ("Stage 1", "\u5b66\u79d1\u5206\u7c7b", "\u8bba\u6587\u2192\u4e3b/\u526f\u5b66\u79d1"),
    ("Stage 2a", "\u6982\u5ff5\u62bd\u53d6", "\u8bba\u6587\u2192\u591a\u5b66\u79d1\u6982\u5ff5"),
    ("Stage 2b", "\u5173\u7cfb\u62bd\u53d6", "\u6982\u5ff5\u2192\u5b66\u79d1\u95f4\u5173\u7cfb"),
    ("Stage 2c", "\u67e5\u8be2\u751f\u6210", "\u6982\u5ff5+\u5173\u7cfb\u2192L1/L2/L3"),
    ("Stage 2d", "\u5047\u8bbe\u8def\u5f84", "\u67e5\u8be2+KG\u2192\u5047\u8bbe\u8def\u5f84"),
    ("Stage 3", "GT \u6784\u5efa", "\u62bd\u53d6\u7ed3\u679c\u2192Ground Truth"),
    ("Stage 4", "P1-P5 \u8bc4\u4f30", "\u63d0\u793a\u6a21\u677f\u2192\u591a\u7ef4\u5ea6\u8bc4\u5206"),
]

for i, (stage, name, desc) in enumerate(stages):
    col = i % 4
    row = i // 4
    x = 0.8 + col * 3.1
    y = 1.5 + row * 2.8
    add_card(slide, x, y, 2.8, 2.2, f"{stage}: {name}", desc, BLUE, LIGHT_BG)

# ══════════════════════════════════════════════════════════════════════
# Slide 3: Example Paper
# ══════════════════════════════════════════════════════════════════════
slide = add_slide("2. \u793a\u4f8b\u8bba\u6587 (Example Paper)")

add_card(slide, 0.8, 1.5, 11.5, 5.0,
         "INPUT: \u8bba\u6587\u4fe1\u606f",
         "",
         BLUE, LIGHT_BG)

add_text_box(slide, 1.0, 2.1, 11, 0.5,
             f"Title: {title}",
             font_size=14, bold=True)

abstract = paper0.get("abstract", "")
add_text_box(slide, 1.0, 2.7, 11, 3.0,
             f"Abstract: {abstract[:400]}...",
             font_size=11, color=GRAY)

add_text_box(slide, 1.0, 5.8, 5, 0.4,
             f"Primary: {primary}  |  Secondary: {secondary}",
             font_size=12, bold=True, color=GREEN)

# ══════════════════════════════════════════════════════════════════════
# Slide 4: Discipline Classification
# ══════════════════════════════════════════════════════════════════════
slide = add_slide("3. Stage 1: \u5b66\u79d1\u5206\u7c7b (Discipline Classification)")

add_card(slide, 0.8, 1.5, 5.2, 2.5,
         "INPUT",
         f"Title: {title[:60]}...\nAbstract: {abstract[:100]}...",
         BLUE, LIGHT_BG)

add_card(slide, 6.5, 1.5, 5.5, 2.5,
         "OUTPUT",
         f"\u4e3b\u5b66\u79d1: {primary}\n\u526f\u5b66\u79d1: {secondary}\n\u2192 \u8bc6\u522b\u4e3a\u8de8\u5b66\u79d1\u8bba\u6587\uff0c\u6d89\u53ca {len(paper0.get('secondary_list', []))} \u4e2a\u9886\u57df",
         GREEN, LIGHT_BG)

# Arrow between cards
arrow = slide.shapes.add_shape(
    13, Inches(6.1), Inches(2.5), Inches(0.3), Inches(0.3))
arrow.fill.solid()
arrow.fill.fore_color.rgb = ORANGE
arrow.line.fill.background()

# ══════════════════════════════════════════════════════════════════════
# Slide 5: Concept Extraction
# ══════════════════════════════════════════════════════════════════════
slide = add_slide("4. Stage 2a: \u6982\u5ff5\u62bd\u53d6 (Concept Extraction)")

concepts = paper0.get("concepts", {})
total_concepts = sum(len(v) for v in concepts.values())

add_text_box(slide, 0.8, 1.4, 11.5, 0.5,
             f"\u5171\u62bd\u53d6 {total_concepts} \u4e2a\u6982\u5ff5\uff0c\u8986\u76d6 {len(concepts)} \u4e2a\u5b66\u79d1\u9886\u57df",
             font_size=14, bold=True, color=GREEN)

concept_rows = []
for disc, clist in concepts.items():
    for c in clist[:4]:
        concept_rows.append([
            c.get("term", "")[:40],
            c.get("normalized", ""),
            disc,
            f"{c.get('confidence', 0):.2f}",
        ])

add_table_slide(slide, 0.8, 2.0, 11.5, 4.5,
                ["Term", "Normalized", "Discipline", "Confidence"],
                concept_rows[:12])

# ══════════════════════════════════════════════════════════════════════
# Slide 6: Relation Extraction
# ══════════════════════════════════════════════════════════════════════
slide = add_slide("5. Stage 2b: \u5173\u7cfb\u62bd\u53d6 (Relation Extraction)")

relations = paper0.get("relations", [])

add_text_box(slide, 0.8, 1.4, 11.5, 0.5,
             f"\u5171\u62bd\u53d6 {len(relations)} \u6761\u5173\u7cfb",
             font_size=14, bold=True, color=GREEN)

rel_rows = []
for r in relations:
    rel_rows.append([
        r.get("head", "")[:35],
        r.get("relation_type", ""),
        r.get("tail", "")[:35],
        f"{r.get('confidence', 0):.2f}",
    ])

add_table_slide(slide, 0.8, 2.0, 11.5, 3.5,
                ["Head", "Relation Type", "Tail", "Conf."],
                rel_rows)

# ══════════════════════════════════════════════════════════════════════
# Slide 7: Query Generation
# ══════════════════════════════════════════════════════════════════════
slide = add_slide("6. Stage 2c: \u67e5\u8be2\u751f\u6210 (Query Generation)")

queries = paper0.get("queries", {})
l1_q = queries.get("L1", "")
l2_q = queries.get("L2", [])
l3_q = queries.get("L3", [])

add_card(slide, 0.8, 1.5, 11.5, 1.2,
         "L1 \u5b8f\u89c2\u67e5\u8be2",
         l1_q if isinstance(l1_q, str) else str(l1_q),
         BLUE, LIGHT_BG)

l2_text = "\n".join(l2_q) if isinstance(l2_q, list) else str(l2_q)
add_card(slide, 0.8, 3.0, 11.5, 1.5,
         "L2 \u4e2d\u89c2\u67e5\u8be2",
         l2_text,
         GREEN, LIGHT_BG)

l3_text = "\n".join(l3_q) if isinstance(l3_q, list) else str(l3_q)
add_card(slide, 0.8, 4.8, 11.5, 1.5,
         "L3 \u5fae\u89c2\u67e5\u8be2",
         l3_text,
         ORANGE, LIGHT_BG)

# ══════════════════════════════════════════════════════════════════════
# Slide 8: Hypothesis Paths
# ══════════════════════════════════════════════════════════════════════
slide = add_slide("7. Stage 2d: \u5047\u8bbe\u8def\u5f84\u751f\u6210 (Hypothesis Paths)")

hyp_paths = paper0.get("hypothesis_paths", {})
total_paths = sum(len(v) if isinstance(v, list) else 0 for v in hyp_paths.values())

add_text_box(slide, 0.8, 1.4, 11.5, 0.5,
             f"\u5171\u751f\u6210 {total_paths} \u6761\u5047\u8bbe\u8def\u5f84 (L1: {len(hyp_paths.get('L1', []))}, L2: {len(hyp_paths.get('L2', []))}, L3: {len(hyp_paths.get('L3', []))})",
             font_size=14, bold=True, color=GREEN)

# Show L1 first path
l1_paths = hyp_paths.get("L1", [])
if l1_paths:
    path1 = l1_paths[0]
    steps = path1.get("steps", path1.get("path", []))
    if isinstance(steps, list):
        path_text = ""
        for si, step in enumerate(steps):
            if isinstance(step, dict):
                head = step.get("head", "")
                tail = step.get("tail", "")
                rel = step.get("relation", "")[:60]
                path_text += f"Step {si+1}: {head} \u2192 {tail}\n  [{rel}]\n"
        add_card(slide, 0.8, 2.1, 5.5, 4.0,
                 "L1 Path 1",
                 path_text,
                 BLUE, LIGHT_BG)

# Show L2 first path
l2_paths = hyp_paths.get("L2", [])
if l2_paths:
    path1 = l2_paths[0]
    steps = path1.get("steps", path1.get("path", []))
    if isinstance(steps, list):
        path_text = ""
        for si, step in enumerate(steps):
            if isinstance(step, dict):
                head = step.get("head", "")
                tail = step.get("tail", "")
                rel = step.get("relation", "")[:60]
                path_text += f"Step {si+1}: {head} \u2192 {tail}\n  [{rel}]\n"
        add_card(slide, 6.6, 2.1, 5.5, 4.0,
                 "L2 Path 1",
                 path_text,
                 GREEN, LIGHT_BG)

# ══════════════════════════════════════════════════════════════════════
# Slide 9: KG Metrics
# ══════════════════════════════════════════════════════════════════════
slide = add_slide("8. \u77e5\u8bc6\u56fe\u8c31\u6784\u5efa\u4e0e\u6307\u6807 (KG Construction & Metrics)")

add_table_slide(slide, 0.8, 1.5, 5, 2.5,
                ["\u6307\u6807", "\u503c"],
                [
                    ["\u5904\u7406\u8bba\u6587\u6570", str(n_papers)],
                    ["\u5e73\u5747\u6982\u5ff5/\u7bc7", f"{avg_c:.1f}"],
                    ["\u5e73\u5747\u5173\u7cfb/\u7bc7", f"{avg_r:.1f}"],
                    ["\u603b\u5173\u7cfb\u6570", str(sum(all_relations_counts))],
                ])

if kg0:
    scores = kg0.get("scores", {})
    kg_rows = []
    for level in ["L1", "L2", "L3"]:
        kg_rows.append([
            level,
            f"{scores.get(f'{level}_bridging', 0):.2f}",
            f"{scores.get(f'{level}_chain_coherence', 0):.3f}",
            f"{scores.get(f'{level}_innovation', 0):.1f}",
            f"{scores.get(f'{level}_feasibility', 0):.1f}",
            f"{scores.get(f'{level}_entity_coverage', 0):.2f}",
        ])
    add_table_slide(slide, 6.2, 1.5, 6.3, 2.5,
                    ["Level", "Bridging", "Coherence", "Innovation", "Feasibility", "Coverage"],
                    kg_rows)

# Depth metrics
add_text_box(slide, 0.8, 4.5, 11.5, 0.5,
             "\u6df1\u5ea6\u5206\u6790 (Depth Analysis)", font_size=16, bold=True, color=BLUE)

if kg0:
    scores = kg0.get("scores", {})
    add_table_slide(slide, 0.8, 5.2, 11.5, 1.8,
                    ["\u6307\u6807", "L2 Concept Exp.", "L3 Concept Exp.", "L2 Anchoring", "L3 Anchoring", "Depth Quality"],
                    [["\u503c",
                      f"{scores.get('depth_l2_concept_expansion', 0)}",
                      f"{scores.get('depth_l3_concept_expansion', 0)}",
                      f"{scores.get('depth_l2_anchoring', 0)}",
                      f"{scores.get('depth_l3_anchoring', 0)}",
                      f"{scores.get('depth_depth_quality', 0)}"]])

# ══════════════════════════════════════════════════════════════════════
# Slide 10: GT Construction
# ══════════════════════════════════════════════════════════════════════
slide = add_slide("9. Benchmark GT \u6784\u5efa (Ground Truth Construction)")

avg_gt_t = sum(all_gt_terms) / max(len(all_gt_terms), 1)
avg_gt_r = sum(all_gt_rels) / max(len(all_gt_rels), 1)

add_table_slide(slide, 0.8, 1.5, 5, 2.5,
                ["\u6307\u6807", "\u503c"],
                [
                    ["GT \u6761\u76ee\u6570", str(n_gt)],
                    ["\u5e73\u5747\u672f\u8bed/\u6761", f"{avg_gt_t:.1f}"],
                    ["\u5e73\u5747\u5173\u7cfb/\u6761", f"{avg_gt_r:.1f}"],
                    ["\u603b\u672f\u8bed\u6570", str(sum(all_gt_terms))],
                    ["\u603b\u5173\u7cfb\u6570", str(sum(all_gt_rels))],
                ])

# Example GT terms
gt_data = gt0.get("ground_truth", {})
gt_terms = gt_data.get("terms", [])
gt_term_rows = []
for t in gt_terms[:6]:
    gt_term_rows.append([
        t.get("term", "")[:30],
        t.get("normalized", ""),
        t.get("discipline", ""),
        f"{t.get('confidence', 0):.2f}",
    ])
add_table_slide(slide, 6.2, 1.5, 6.3, 3.5,
                ["Term", "Normalized", "Discipline", "Conf."],
                gt_term_rows)

# GT relations
gt_rels = gt_data.get("relations", [])
gt_rel_rows = []
for r in gt_rels[:4]:
    gt_rel_rows.append([
        r.get("head", "")[:20],
        r.get("relation_type", ""),
        r.get("tail", "")[:20],
        f"{r.get('confidence', 0):.1f}",
    ])
add_table_slide(slide, 0.8, 5.0, 11.5, 2.0,
                ["Head", "Type", "Tail", "Conf."],
                gt_rel_rows)

# ══════════════════════════════════════════════════════════════════════
# Slide 11: P1-P5 Evaluation
# ══════════════════════════════════════════════════════════════════════
slide = add_slide("10. P1-P5 \u8bc4\u4f30\u6846\u67b6 (Progressive Prompt Evaluation)")

p_levels = ["P1", "P2", "P3", "P4", "P5"]
eval_rows = []
for pl in p_levels:
    d = p1p5.get(pl, {})
    novelty = d.get('judge_novelty')
    feasibility = d.get('judge_feasibility')
    relevance = d.get('judge_relevance')
    cross_disc = d.get('judge_cross_disciplinary')
    eval_rows.append([
        pl,
        f"{d.get('text_bertscore_f1', 0):.4f}",
        f"{d.get('text_rouge1', 0):.4f}",
        f"{novelty}" if novelty is not None else "N/A",
        f"{feasibility}" if feasibility is not None else "N/A",
        f"{relevance}" if relevance is not None else "N/A",
        f"{cross_disc}" if cross_disc is not None else "N/A",
        f"{d.get('elapsed_seconds', 0):.1f}s",
    ])

add_table_slide(slide, 0.8, 1.5, 11.5, 3.5,
                ["Level", "BERTScore", "ROUGE-1", "Novelty", "Feasibility", "Relevance", "Cross-Disc", "Time"],
                eval_rows)

# Key findings
add_bullet_list(slide, 0.8, 5.3, 11.5, 2.0, [
    "P4 \u6587\u672c\u76f8\u4f3c\u5ea6\u6700\u4f18 (BERTScore F1 = 0.7472)",
    "P1 \u65b0\u9896\u6027\u6700\u9ad8 (9.0)\uff0cP5 \u76f8\u5173\u6027\u6700\u9ad8 (10.0)",
    "P5 \u751f\u6210\u6700\u5feb\uff0c\u5e73\u5747 6.2 \u6761\u5047\u8bbe",
    "\u4fe1\u606f\u5145\u5206\u5ea6\u4e0e\u65b0\u9896\u6027\u8d1f\u76f8\u5173\uff0c\u4e0e\u76f8\u5173\u6027/\u53ef\u884c\u6027\u6b63\u76f8\u5173",
], font_size=13)

# ══════════════════════════════════════════════════════════════════════
# Slide 12: Depth Comparison L1/L2/L3
# ══════════════════════════════════════════════════════════════════════
slide = add_slide("11. \u6df1\u5ea6\u5206\u6790 (L1 vs L2 vs L3)")

if kg0:
    scores = kg0.get("scores", {})
    depth_rows = []
    for metric in ["innovation", "feasibility", "scientificity", "testability", "chain_coherence", "entity_coverage"]:
        depth_rows.append([
            metric.replace("_", " ").title(),
            f"{scores.get(f'L1_{metric}', 0):.2f}" if isinstance(scores.get(f'L1_{metric}'), float) else str(scores.get(f'L1_{metric}', 'N/A')),
            f"{scores.get(f'L2_{metric}', 0):.2f}" if isinstance(scores.get(f'L2_{metric}'), float) else str(scores.get(f'L2_{metric}', 'N/A')),
            f"{scores.get(f'L3_{metric}', 0):.2f}" if isinstance(scores.get(f'L3_{metric}'), float) else str(scores.get(f'L3_{metric}', 'N/A')),
        ])
    add_table_slide(slide, 0.8, 1.5, 11.5, 4.0,
                    ["\u6307\u6807", "L1 (\u5b8f\u89c2)", "L2 (\u4e2d\u89c2)", "L3 (\u5fae\u89c2)"],
                    depth_rows)

add_bullet_list(slide, 0.8, 5.8, 11.5, 1.5, [
    "L1: \u5b66\u79d1\u95f4\u5b8f\u89c2\u8054\u7cfb\uff0c\u9ad8\u521b\u65b0\u6027\uff0c\u5c11\u91cf\u6982\u5ff5",
    "L2: \u7279\u5b9a\u673a\u5236\u7ea7\u5173\u8054\uff0c\u6982\u5ff5\u6269\u5c55 25%",
    "L3: \u7ec6\u7c92\u5ea6\u5b9e\u9a8c\u7ea7\u5173\u8054\uff0c\u6982\u5ff5\u6269\u5c55 62.5%",
], font_size=13)

# ══════════════════════════════════════════════════════════════════════
# Slide 13: Summary
# ══════════════════════════════════════════════════════════════════════
slide = add_slide("12. \u6570\u636e\u96c6\u7edf\u8ba1\u6982\u89c8 (Dataset Summary)")

add_table_slide(slide, 0.8, 1.5, 5, 4.5,
                ["\u7ef4\u5ea6", "\u7edf\u8ba1\u503c"],
                [
                    ["\u8bba\u6587\u6570\u91cf", str(n_papers)],
                    ["GT \u6761\u76ee\u6570", str(n_gt)],
                    ["\u5e73\u5747\u6982\u5ff5/\u7bc7", f"{avg_c:.1f}"],
                    ["\u5e73\u5747\u5173\u7cfb/\u7bc7", f"{avg_r:.1f}"],
                    ["\u5e73\u5747 GT \u672f\u8bed/\u6761", f"{avg_gt_t:.1f}"],
                    ["\u5e73\u5747 GT \u5173\u7cfb/\u6761", f"{avg_gt_r:.1f}"],
                    ["\u603b\u672f\u8bed\u6570", str(sum(all_gt_terms))],
                    ["\u603b\u5173\u7cfb\u6570", str(sum(all_gt_rels))],
                ])

add_bullet_list(slide, 6.2, 1.5, 5.8, 5.0, [
    f"\u5904\u7406 {n_papers} \u7bc7 Nature Communications \u8bba\u6587",
    f"\u6784\u5efa {n_gt} \u6761 Evidence-Grounded GT",
    f"\u5e73\u5747\u6bcf\u7bc7\u62bd\u53d6 {avg_c:.0f} \u4e2a\u8de8\u5b66\u79d1\u6982\u5ff5",
    f"\u5e73\u5747\u6bcf\u7bc7\u62bd\u53d6 {avg_r:.0f} \u6761\u5b66\u79d1\u95f4\u5173\u7cfb",
    "P1-P5 \u6e10\u8fdb\u5f0f\u8bc4\u4f30\u6846\u67b6\u591a\u7ef4\u5ea6\u5206\u6790",
    "\u4e09\u5c42\u6b21\u5047\u8bbe\u8def\u5f84 (L1/L2/L3) \u6df1\u5ea6\u5206\u6790",
    "12 \u79cd\u5173\u7cfb\u7c7b\u578b\u8986\u76d6\u8de8\u5b66\u79d1\u7814\u7a76\u6a21\u5f0f",
], font_size=14)

# ── Save ──
output_path = os.path.join(BASE, "CrossDisc_Bench_pipeline_walkthrough.pptx")
prs.save(output_path)
print(f"PPT saved to: {output_path}")
