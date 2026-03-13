"""生成 Benchmark 项目运行流程详解 PPT."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import json

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 颜色常量 ──
BLUE = RGBColor(0x1A, 0x56, 0xDB)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
CODE_BG = RGBColor(0xF2, 0xF2, 0xF2)
ACCENT = RGBColor(0x0D, 0x92, 0x76)

# ── 辅助函数 ──

def add_slide(title_text, layout_idx=1):
    """添加幻灯片并设置标题."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    # 清除默认占位符
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)
    # 添加标题
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.font.name = "微软雅黑"
    # 底部分隔线
    from pptx.oxml.ns import qn
    line = slide.shapes.add_shape(
        1, Inches(0.8), Inches(1.15), Inches(11.5), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()
    return slide


def add_text_box(slide, left, top, width, height, text, font_size=16,
                 bold=False, color=DARK, font_name="微软雅黑", align=PP_ALIGN.LEFT):
    """添加文本框."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=15, spacing=Pt(6)):
    """添加项目符号列表."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK
        p.font.name = "微软雅黑"
        p.space_after = spacing
    return tf


def add_code_box(slide, left, top, width, height, text, font_size=11):
    """添加代码框（灰色背景）."""
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(8)
    tf.margin_right = Pt(12)
    tf.margin_bottom = Pt(8)
    for i, line in enumerate(text.strip().split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK
        p.font.name = "Consolas"
        p.space_after = Pt(2)
    return shape


def add_table_slide(slide, left, top, width, height, headers, rows, col_widths=None):
    """添加表格."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table

    # 表头
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "微软雅黑"
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE

    # 数据行
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = DARK
                p.font.name = "微软雅黑"
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    return table


def add_subtitle(slide, left, top, text, font_size=20):
    """添加子标题."""
    add_text_box(slide, left, top, 11, 0.5, text, font_size=font_size, bold=True, color=ACCENT)


# ══════════════════════════════════════════════════════════════════════
#  Slide 1: 封面
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
# 背景
bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
bg.line.fill.background()

add_text_box(slide, 1.5, 1.8, 10, 1.2,
             "跨学科知识抽取与基准评测系统",
             font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 3.0, 10, 0.8,
             "Cross-Disciplinary Knowledge Extraction & Benchmarking",
             font_size=20, color=RGBColor(0x88, 0xAA, 0xDD), align=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 4.2, 10, 0.6,
             "项目运行流程与输入输出详解  ·  v0.2.0",
             font_size=22, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════
#  Slide 2: 目录
# ══════════════════════════════════════════════════════════════════════
slide = add_slide("目录")
items = [
    "一、项目概述与核心能力",
    "二、端到端运行总览（6 阶段流程图）",
    "三、Phase 0：数据准备（输入格式与 CLI 命令）",
    "四、Phase 1：学科分类（MSC 层级分类体系）",
    "五、Phase 2：多学科过滤",
    "六、Phase 3：三阶段知识抽取（核心）",
    "    Stage 1a/1b：概念抽取 + 关系抽取",
    "    Stage 2：查询生成（三级层次化）",
    "    Stage 3：假设生成（L1/L2/L3 链式推理）",
    "七、Phase 4：图谱构建与指标计算（15+ 指标）",
    "八、Phase 5：Ground Truth 构建",
    "九、Phase 6：基准评测（多维度评分）",
    "十、CLI 命令参考与配置",
]
add_bullet_list(slide, 1.0, 1.5, 11, 5.5, items, font_size=17, spacing=Pt(8))

# ══════════════════════════════════════════════════════════════════════
#  Slide 3: 项目概述
# ══════════════════════════════════════════════════════════════════════
slide = add_slide("一、项目概述")
add_text_box(slide, 0.8, 1.5, 11.5, 0.8,
             "基于多阶段大语言模型（LLM）Pipeline 的跨学科知识抽取与基准评测系统，"
             "从学术论文中自动抽取跨学科概念、关系，生成科学假设，并通过多维度指标进行评估。",
             font_size=16)

add_subtitle(slide, 0.8, 2.5, "核心能力")
items = [
    "学科分类：基于 MSC 层级分类体系，自动识别论文的主学科与辅学科",
    "概念抽取：从论文标题、摘要、引言中抽取专业术语，按学科分组",
    "关系抽取：识别跨学科概念间的语义关系（11 种标准关系类型）",
    "查询生成：生成三级层次化查询（宏观 → 学科维度 → 细粒度）",
    "假设生成：生成三级知识路径（每条路径 3 步链式推理）",
    "知识图谱：基于 NetworkX 构建概念图谱并计算 15+ 图谱指标",
    "基准评测：LLM-as-Judge 主观评分 + 客观图谱指标的多维度评估",
]
add_bullet_list(slide, 0.8, 3.0, 11.5, 4, items, font_size=15)

# ══════════════════════════════════════════════════════════════════════
#  Slide 4: 技术栈
# ══════════════════════════════════════════════════════════════════════
slide = add_slide("技术栈")
add_table_slide(slide, 0.8, 1.5, 11.5, 4.5,
    ["组件", "技术选型", "用途"],
    [
        ["LLM 调用", "OpenAI 兼容 API (qwen3-235b-a22b)", "概念/关系/假设抽取 + 评估"],
        ["数据模型", "Pydantic v2", "严格 JSON Schema 验证"],
        ["图谱构建", "NetworkX", "概念图谱 + 拓扑指标"],
        ["重试机制", "tenacity", "指数退避重试（最多 5 次）"],
        ["分类器", "LangChain + OpenAI", "层级学科分类"],
        ["PDF 处理", "pdfminer / pdfplumber", "论文引言提取"],
        ["语义相似度", "sentence-transformers", "嵌入桥接 + 链式连贯性"],
        ["配置管理", "PyYAML + frozen dataclass", "线程安全配置"],
    ],
    col_widths=[2.5, 4.5, 4.5],
)

# ══════════════════════════════════════════════════════════════════════
#  Slide 5: 端到端流程图
# ══════════════════════════════════════════════════════════════════════
slide = add_slide("二、端到端运行总览")
add_code_box(slide, 0.8, 1.5, 11.5, 4.8, """
crossdisc-pipeline full -i papers.jsonl -o results.jsonl

Phase 0: 数据准备 (输入论文 JSONL)
    ↓
Phase 1: 学科分类 (classifier/) → 主学科 + 辅学科
    ↓
Phase 2: 多学科过滤 → 仅保留跨学科论文
    ↓
Phase 3: 三阶段知识抽取 (extractor_multi_stage.py)
    ├─ Stage 1a: 概念抽取 (LLM Round 1 + Round 2 补充)
    ├─ Stage 1b: 关系抽取 (11 种标准关系类型)
    ├─ Stage 2:  查询生成 (三级层次化)
    └─ Stage 3:  假设生成 (L1/L2/L3 链式推理)
    ↓
Phase 4: 图谱构建 + 指标计算 (graph_builder.py) → 15+ 指标
    ↓
Phase 5: Ground Truth 构建 (gt_builder.py)
    ↓
Phase 6: 基准评测 (evaluate_benchmark.py) → 多维度评分
""", font_size=14)

# ══════════════════════════════════════════════════════════════════════
#  Slide 6: LLM 调用统计
# ══════════════════════════════════════════════════════════════════════
slide = add_slide("LLM 调用次数统计（单篇论文）")
add_table_slide(slide, 1.5, 1.5, 10, 4,
    ["阶段", "LLM 调用次数", "说明"],
    [
        ["Phase 1 分类", "2-4 次", "每层分类体系一次（通常 2-3 层）"],
        ["Stage 1a 概念", "2 次", "Round 1 + Round 2 补充"],
        ["Stage 1b 关系", "1 次", "基于概念列表抽取关系"],
        ["Stage 2 查询", "1 次", "生成三级查询"],
        ["Stage 3 假设", "3 次", "L1 + L2 + L3 各一次"],
        ["总计", "~9-11 次", "单篇论文的完整处理"],
    ],
    col_widths=[3, 2.5, 4.5],
)

# ══════════════════════════════════════════════════════════════════════
#  Slide 7: Phase 0 数据准备
# ══════════════════════════════════════════════════════════════════════
slide = add_slide("三、Phase 0：数据准备")
add_subtitle(slide, 0.8, 1.5, "CLI 命令")
add_code_box(slide, 0.8, 2.0, 11.5, 0.7,
    "crossdisc-pipeline full -i papers.jsonl -o results.jsonl --config configs/default.yaml",
    font_size=13)

add_subtitle(slide, 0.8, 3.0, "输入格式（JSONL，每行一篇论文）")
add_code_box(slide, 0.8, 3.5, 11.5, 1.5, """{
  "title": "Refining centromedian nucleus stimulation for generalized epilepsy...",
  "abstract": "Epilepsy affects 65 million people worldwide...",
  "pdf_url": "https://www.nature.com/articles/s41467-025-60183-9.pdf"
}""", font_size=12)

add_subtitle(slide, 0.8, 5.3, "支持格式与字段")
items = [
    "支持 .json / .jsonl / .csv 三种格式",
    "字段名兼容中英文变体（标题/title、摘要/abstract）",
    "必需字段：title、abstract | 可选字段：pdf_url、primary、secondary_list",
]
add_bullet_list(slide, 0.8, 5.8, 11.5, 1.5, items, font_size=14)

# ══════════════════════════════════════════════════════════════════════
#  Slide 8: Phase 1 学科分类
# ══════════════════════════════════════════════════════════════════════
slide = add_slide("四、Phase 1：学科分类")
add_subtitle(slide, 0.8, 1.5, "分类流程（逐层遍历 MSC 分类体系）")
items = [
    "Level 0：从顶层学科中选择（数学、物理学、临床医学、基础医学...）",
    "Level 1：在选中的 L0 学科下选择子类",
    "Level 2：继续细分（如有更深层级）",
    "每层 LLM 调用：构建 Prompt → 解析方括号列表输出 → 校验选项合法性",
    "多学科判定：len(partial_paths) > 1 时标记为跨学科论文",
]
add_bullet_list(slide, 0.8, 2.0, 5.5, 3.5, items, font_size=14)

add_subtitle(slide, 6.8, 1.5, "输入 → 输出示例")
add_code_box(slide, 6.8, 2.0, 5.8, 1.2, """输入:
{"title": "Refining centromedian...",
 "abstract": "Epilepsy affects..."}""", font_size=11)

add_code_box(slide, 6.8, 3.4, 5.8, 1.0, """LLM (Level 0):
Prompt: "请从以下学科中选择..."
输出: "[临床医学, 基础医学, 工程与技术科学基础学科]" """, font_size=11)

add_code_box(slide, 6.8, 4.6, 5.8, 2.2, """输出:
{
  "primary": "临床医学",
  "secondary_list": [
    "基础医学",
    "工程与技术科学基础学科"
  ],
  "is_multidisciplinary": true
}""", font_size=11)

# ══════════════════════════════════════════════════════════════════════
#  Slide 9: Phase 2 过滤
# ══════════════════════════════════════════════════════════════════════
slide = add_slide("五、Phase 2：多学科过滤")
items = [
    "检查分类结果中 secondary_list 是否非空",
    "仅保留跨学科论文（secondary_list 非空）进入后续抽取阶段",
    "单学科论文被过滤掉，不进入知识抽取流程",
    "过滤目的：只有涉及多个学科的论文才具备跨学科知识桥接的分析价值",
]
add_bullet_list(slide, 0.8, 1.8, 11.5, 3, items, font_size=18, spacing=Pt(16))

add_code_box(slide, 2, 4.5, 9, 1.8, """# pipeline.py 中 full 子命令逻辑
multi_papers = [p for p in classified_papers if p["secondary_list"]]
print(f"共 {len(classified_papers)} 篇论文，其中 {len(multi_papers)} 篇为跨学科论文")
# 仅 multi_papers 进入 Phase 3 抽取阶段""", font_size=14)

# ══════════════════════════════════════════════════════════════════════
#  Slide 10: Stage 1a 概念抽取
# ══════════════════════════════════════════════════════════════════════
slide = add_slide("六、Phase 3 — Stage 1a：概念抽取")
add_text_box(slide, 0.8, 1.4, 11.5, 0.5,
    "调用链：build_concepts_messages() → LLM → parse_concepts_output()", font_size=14, color=GRAY)

add_subtitle(slide, 0.8, 1.9, "输入变量")
add_code_box(slide, 0.8, 2.4, 5.5, 1.8, """title = "Refining centromedian..."
abstract = "Epilepsy affects 65 million..."
introduction = "[从 PDF 提取的引言]"
primary = "临床医学"
secondary_list = "基础医学, 工程与技术科学基础学科" """, font_size=11)

add_subtitle(slide, 6.8, 1.9, "输出示例")
add_code_box(slide, 6.8, 2.4, 5.8, 4.5, """{
  "meta": {"title": "...", "primary": "临床医学",
           "secondary_list": ["基础医学", "工程与技术科学基础学科"]},
  "概念": {
    "主学科": [
      {"term": "癫痫", "normalized": "Epilepsy",
       "evidence": "Epilepsy affects 65 million...",
       "source": "abstract", "confidence": 1.0},
      {"term": "脑深部刺激", "normalized": "Deep brain stimulation",
       "confidence": 1.0}
    ],
    "辅学科": {
      "基础医学": [
        {"term": "神经学", "normalized": "Neurology", "confidence": 0.9},
        {"term": "电生理学", "normalized": "Electrophysiology", "confidence": 1.0}
      ],
      "工程与技术科学基础学科": [
        {"term": "信号处理", "normalized": "Signal processing", "confidence": 0.8}
      ]
    }
  }
}""", font_size=10)

add_subtitle(slide, 0.8, 4.5, "Round 2 补充抽取")
items = [
    "将 Round 1 概念列表传入 LLM，要求补充遗漏术语",
    "通过 _merge_concepts() 去重合并",
    "通过 _ground_and_filter_concepts() 过滤学科名/通用词",
]
add_bullet_list(slide, 0.8, 5.0, 5.5, 2, items, font_size=13)

# ══════════════════════════════════════════════════════════════════════
#  Slide 11: Stage 1b 关系抽取
# ══════════════════════════════════════════════════════════════════════
slide = add_slide("Phase 3 — Stage 1b：关系抽取")
add_text_box(slide, 0.8, 1.4, 11.5, 0.5,
    "调用链：build_relations_messages() → LLM → parse_relations_output()", font_size=14, color=GRAY)

add_subtitle(slide, 0.8, 1.9, "输出示例")
add_code_box(slide, 0.8, 2.4, 6.5, 4.0, """{
  "跨学科关系": [
    {
      "head": "Deep brain stimulation",
      "relation": "offers a promising alternative treatment for",
      "relation_type": "method_applied_to",
      "tail": "Epilepsy",
      "direction": "->",
      "assumptions": ["Drug-resistant epilepsy is present"],
      "evidence": "Centromedian nucleus neurostimulation offers...",
      "confidence": 0.9
    },
    {
      "head": "Electrophysiology",
      "relation": "guides the optimization of",
      "relation_type": "improves_metric",
      "tail": "Deep brain stimulation",
      "confidence": 0.9
    }
  ]
}""", font_size=11)

add_subtitle(slide, 7.8, 1.9, "11 种标准关系类型")
rel_items = [
    "method_applied_to（方法应用于）",
    "maps_to（映射到）",
    "constrains（约束）",
    "improves_metric（改善指标）",
    "corresponds_to（对应于）",
    "inferred_from（推断自）",
    "assumes（假设）",
    "extends（扩展）",
    "generalizes（泛化）",
    "driven_by（驱动于）",
    "depends_on（依赖于）",
]
add_bullet_list(slide, 7.8, 2.4, 4.8, 4.5, rel_items, font_size=13, spacing=Pt(4))

# ══════════════════════════════════════════════════════════════════════
#  Slide 12: Stage 2 查询生成
# ══════════════════════════════════════════════════════════════════════
slide = add_slide("Phase 3 — Stage 2：查询生成")
add_text_box(slide, 0.8, 1.4, 11.5, 0.5,
    "调用链：build_query_messages() → LLM → parse_query_output()", font_size=14, color=GRAY)

add_subtitle(slide, 0.8, 1.9, "输出示例")
add_code_box(slide, 0.8, 2.4, 7.0, 4.2, """{
  "按辅助学科分类": {
    "基础医学": {
      "概念": ["神经学", "神经生理学", "神经解剖学", "电生理学", "神经回路"],
      "关系": [0, 1, 2, 3, 4, 5],
      "rationale": "基础医学提供神经机制、结构和电生理基础"
    },
    "工程与技术科学基础学科": {
      "概念": ["影像引导技术", "信号处理", "神经工程"],
      "rationale": "工程学科支撑高精度定位与信号分析"
    }
  },
  "查询": {
    "一级": "如何通过多学科技术优化脑深部刺激治疗癫痫？",
    "二级": ["结合基础医学中的神经生理学与神经解剖学",
            "利用工程科学的影像引导与信号处理技术"],
    "三级": ["如何协同运用神经解剖学、信号处理实现精准靶向",
            "通过多模态成像引导验证脑深部刺激的机制与效果"]
  }
}""", font_size=11)

add_subtitle(slide, 8.3, 1.9, "查询层级含义")
add_table_slide(slide, 8.3, 2.5, 4.3, 2.5,
    ["层级", "含义"],
    [
        ["一级 L1", "宏观/论文级核心问题"],
        ["二级 L2", "按辅学科维度拆分的子问题"],
        ["三级 L3", "具体操作/方法的细粒度问题"],
    ],
    col_widths=[1.3, 3.0],
)

# ══════════════════════════════════════════════════════════════════════
#  Slide 13: Stage 3 假设生成
# ══════════════════════════════════════════════════════════════════════
slide = add_slide("Phase 3 — Stage 3：假设生成（L1/L2/L3）")
add_text_box(slide, 0.8, 1.4, 11.5, 0.5,
    "分三次独立调用 LLM，分别生成 L1（宏观）、L2（学科维度）、L3（细粒度）假设",
    font_size=14, color=GRAY)

add_subtitle(slide, 0.8, 1.9, "L1 输出示例（宏观假设 — 每条路径 3 步链式推理）")
add_code_box(slide, 0.8, 2.4, 11.5, 3.5, """{
  "假设": {
    "一级": [
      [
        {"step": 1, "head": "脑深部刺激", "relation": "应用于", "tail": "癫痫治疗",
         "claim": "脑深部刺激是治疗药物难治性癫痫的一种有前景的方法。"},
        {"step": 2, "head": "癫痫治疗", "relation": "面临", "tail": "作用机制不明确",
         "claim": "当前脑深部刺激治疗癫痫的神经调控机制尚不完全清楚。"},
        {"step": 3, "head": "作用机制不明确", "relation": "需要", "tail": "基础医学补位",
         "claim": "需要基础医学通过神经生理学和神经解剖学揭示调控机制。"}
      ]
    ],
    "一级总结": ["脑深部刺激治疗癫痫需要基础医学揭示调控机制以优化刺激靶点"]
  }
}""", font_size=11)

add_subtitle(slide, 0.8, 6.1, "关键约束")
items = [
    "每条路径严格 3 步：step[i].tail == step[i+1].head（相似度 ≥ 0.75）",
    "head/tail 必须来自已抽取的概念列表",
    "最后一步的 claim 必须非空，且为完整可验证的假设陈述",
]
add_bullet_list(slide, 0.8, 6.5, 11.5, 1, items, font_size=13, spacing=Pt(3))

# ══════════════════════════════════════════════════════════════════════
#  Slide 14: Phase 4 图谱构建
# ══════════════════════════════════════════════════════════════════════
slide = add_slide("七、Phase 4：图谱构建与指标计算")
add_text_box(slide, 0.8, 1.4, 11.5, 0.5,
    "入口：graph_builder.py → build_graph_and_metrics()", font_size=14, color=GRAY)

add_subtitle(slide, 0.8, 1.9, "构建流程")
items = [
    "添加节点：从 Stage 1 概念中提取所有术语作为图节点，标记学科归属",
    "添加结构边：从 Stage 1 关系中提取 head→tail 作为图边",
    "添加假设边：从 Stage 3 每个 HypothesisStep 的 head→tail 作为图边",
    "计算指标：15+ 维度的图谱质量指标",
]
add_bullet_list(slide, 0.8, 2.4, 5.5, 2, items, font_size=14)

add_subtitle(slide, 0.8, 4.5, "输出示例")
add_code_box(slide, 0.8, 5.0, 5.5, 2.0, """{
  "graph": {
    "nodes": [
      {"id": "癫痫", "discipline": "临床医学"},
      {"id": "电生理学", "discipline": "基础医学"},
      {"id": "信号处理", "discipline": "工程与技术科学基础学科"}
    ],
    "edges": [
      {"source": "脑深部刺激", "target": "癫痫", "relation": "method_applied_to"}
    ]
  }
}""", font_size=10)

add_subtitle(slide, 6.8, 1.9, "15+ 指标一览")
add_table_slide(slide, 6.8, 2.4, 5.8, 4.5,
    ["指标", "含义"],
    [
        ["path_consistency", "假设路径与原文关系的一致程度"],
        ["coverage", "假设路径对抽取概念的覆盖广度"],
        ["bridging_score", "跨学科知识桥接程度"],
        ["rao_stirling_diversity", "学科多样性/均衡性/差异性"],
        ["embedding_bridging", "首尾概念的语义距离"],
        ["chain_coherence", "推理链条的逻辑连贯程度"],
        ["kg_density", "知识图谱的连接密度"],
        ["kg_modularity", "图谱的社区结构清晰度"],
        ["kg_betweenness", "关键桥接节点的重要性"],
        ["kg_clustering", "局部连接的紧密程度"],
        ["atypical_combination", "概念组合的新颖性"],
    ],
    col_widths=[2.5, 3.3],
)

# ══════════════════════════════════════════════════════════════════════
#  Slide 15: Phase 5 GT 构建
# ══════════════════════════════════════════════════════════════════════
slide = add_slide("八、Phase 5：Ground Truth 构建")
add_text_box(slide, 0.8, 1.4, 11.5, 0.5,
    "入口：benchmark/gt_builder.py → build_ground_truth()", font_size=14, color=GRAY)

add_subtitle(slide, 0.8, 1.9, "三阶段构建流程")

add_text_box(slide, 0.8, 2.5, 3.5, 0.4, "Stage A — 术语抽取", font_size=16, bold=True, color=ACCENT)
items_a = [
    "优先级：已解析概念 > LLM 抽取 > 启发式抽取",
    "启发式：缩写词、科学短语、首字母大写短语",
    "与 MSC 术语库相似度对齐（阈值 0.70）",
]
add_bullet_list(slide, 0.8, 2.9, 3.5, 1.5, items_a, font_size=12, spacing=Pt(3))

add_text_box(slide, 4.8, 2.5, 3.5, 0.4, "Stage B — 关系构建", font_size=16, bold=True, color=ACCENT)
items_b = [
    "原文分句 → 术语共现对",
    "LLM 或启发式分类关系类型",
    "每条关系携带原文证据句",
]
add_bullet_list(slide, 4.8, 2.9, 3.5, 1.5, items_b, font_size=12, spacing=Pt(3))

add_text_box(slide, 8.8, 2.5, 3.5, 0.4, "Stage C — 路径构建", font_size=16, bold=True, color=ACCENT)
items_c = [
    "NetworkX 建有向图",
    "nx.all_simple_paths() 寻路（最长 4 跳）",
    "优先选择跨越多学科的路径",
]
add_bullet_list(slide, 8.8, 2.9, 3.5, 1.5, items_c, font_size=12, spacing=Pt(3))

add_subtitle(slide, 0.8, 4.8, "输出示例")
add_code_box(slide, 0.8, 5.3, 11.5, 1.8, """{
  "id": "Refining centromedian nucleus stimulation...",
  "input": {"title": "...", "primary_discipline": "临床医学",
            "secondary_disciplines": ["基础医学", "工程与技术科学基础学科"]},
  "ground_truth": {
    "terms": [{"term": "Neurological Surgery", "source": "heuristic", "confidence": 0.4}, ...],
    "relations": [{"head": "neural engineering", "tail": "electrophysiology", "relation_type": "depends_on",...}],
    "paths": [...], "concept_graph": {"nodes": [...], "edges": [...]}
  }
}""", font_size=10)

# ══════════════════════════════════════════════════════════════════════
#  Slide 16: Phase 6 基准评测
# ══════════════════════════════════════════════════════════════════════
slide = add_slide("九、Phase 6：基准评测")
add_text_box(slide, 0.8, 1.4, 11.5, 0.5,
    "入口：benchmark/evaluate_benchmark.py", font_size=14, color=GRAY)

add_subtitle(slide, 0.8, 1.9, "评测维度")
add_table_slide(slide, 0.8, 2.4, 7.5, 4.8,
    ["评测维度", "方法", "分值", "说明"],
    [
        ["Innovation", "LLM-as-Judge", "0-10", "假设的新颖程度"],
        ["Feasibility", "LLM-as-Judge", "0-10", "实验可行性"],
        ["Scientificity", "LLM-as-Judge", "0-10", "科学严谨性"],
        ["Consistency", "链式一致性", "0-1", "逻辑一致性"],
        ["Bridging", "概念距离", "0-1", "跨学科桥接"],
        ["Rao-Stirling", "多样性公式", "0-1", "学科多样性"],
        ["Chain Coherence", "SBERT", "0-1", "语义连贯"],
        ["Atypical Comb.", "z-score", "连续", "组合新颖性"],
        ["KG Topology", "图论指标", "连续", "图谱拓扑"],
        ["Concept Coverage", "GT 匹配", "0-1", "术语覆盖"],
        ["Path Alignment", "嵌入相似", "0-1", "路径对齐"],
    ],
    col_widths=[2.2, 1.8, 1.0, 2.5],
)

add_subtitle(slide, 8.8, 1.9, "输出示例")
add_code_box(slide, 8.8, 2.4, 3.8, 4.5, """{
  "id": "-33484234247962345",
  "scores": {
    "L1_innovation": 7.35,
    "L1_feasibility": 7.80,
    "L1_scientificity": 8.70,
    "L1_consistency": 1.0,
    "L1_bridging": 1.0,
    "L2_innovation": 7.40,
    "L2_feasibility": 7.25,
    "L2_scientificity": 6.45,
    "L2_consistency": 0.5,
    "L2_bridging": 1.0,
    "L3_innovation": 8.05,
    "L3_feasibility": 5.25,
    "L3_scientificity": 6.70,
    "L3_consistency": 0.0,
    "L3_bridging": 1.0
  }
}""", font_size=9)

# ══════════════════════════════════════════════════════════════════════
#  Slide 17: CLI 命令参考
# ══════════════════════════════════════════════════════════════════════
slide = add_slide("十、CLI 命令参考与配置")
add_subtitle(slide, 0.8, 1.5, "方式一：端到端 Pipeline（推荐）")
add_code_box(slide, 0.8, 2.0, 5.5, 2.0, """# 完整流水线
crossdisc-pipeline full \\
  -i papers.jsonl \\
  -o results.jsonl \\
  --config configs/default.yaml

# 仅分类
crossdisc-pipeline classify \\
  -i papers.jsonl -o classified.jsonl

# 仅抽取
crossdisc-pipeline extract \\
  -i classified.jsonl -o results.jsonl""", font_size=11)

add_subtitle(slide, 0.8, 4.3, "方式二：独立三阶段抽取")
add_code_box(slide, 0.8, 4.8, 5.5, 1.5, """python run.py batch \\
  --input data.jsonl --output out.jsonl \\
  --num-workers 4 --resume

python run.py single \\
  --title "..." --abstract "..." \\
  --primary "临床医学" --secondary "基础医学,..." """, font_size=11)

add_subtitle(slide, 6.8, 1.5, "关键配置参数")
add_table_slide(slide, 6.8, 2.0, 5.8, 4.5,
    ["参数", "默认值", "说明"],
    [
        ["language_mode", "chinese", "chinese | original"],
        ["temperature_struct", "0.2", "Stage 1 温度"],
        ["temperature_query", "0.2", "Stage 2 温度"],
        ["temperature_hyp", "0.3", "Stage 3 温度"],
        ["seed", "42", "可复现性"],
        ["num_workers", "1", "并行数（1=串行）"],
        ["resume", "true", "断点续传"],
        ["max_tokens_struct", "8192", "Stage 1 max token"],
        ["max_tokens_hyp", "4096", "Stage 3 max token"],
    ],
    col_widths=[2.2, 1.3, 2.3],
)

# ══════════════════════════════════════════════════════════════════════
#  Slide 18: 最终输出结构
# ══════════════════════════════════════════════════════════════════════
slide = add_slide("最终输出数据结构（Extraction）")
add_code_box(slide, 0.8, 1.5, 5.8, 5.5, """{
  "meta": {title, primary, secondary_list},
  "概念": {
    "主学科": [ConceptEntry...],
    "辅学科": {"学科名": [ConceptEntry...]}
  },
  "跨学科关系": [RelationEntry...],
  "按辅助学科分类": {
    "学科名": {
      "概念": [...],
      "关系": [...],
      "rationale": "..."
    }
  },
  "查询": {
    "一级": str,
    "二级": [str],
    "三级": [str]
  },
  "假设": {
    "一级": [[Step,Step,Step]...],
    "一级总结": [str...],
    "二级/三级": ...
  },
  "graph": {"nodes": [...], "edges": [...]},
  "metrics": {15+ 指标}
}""", font_size=12)

add_subtitle(slide, 7.0, 1.5, "数据模型层级")
add_code_box(slide, 7.0, 2.0, 5.5, 4.0, """Extraction（完整抽取结果）
├── meta: {title, primary, secondary_list}
├── 概念 (Concepts)
│   └── ConceptEntry:
│       {term, normalized, std_label,
│        evidence, source, confidence}
├── 跨学科关系: RelationEntry[]
│   └── RelationEntry:
│       {head, tail, relation,
│        relation_type, evidence, confidence}
├── 查询 (Query3Levels): L1/L2/L3
├── 假设 (Hypothesis3Levels)
│   └── HypothesisStep:
│       {step, head, relation, tail, claim}
├── graph (ConceptGraph): nodes + edges
└── metrics (GraphMetrics): 15+ 指标""", font_size=11)

# ══════════════════════════════════════════════════════════════════════
#  Slide 19: 谢谢
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
bg.line.fill.background()

add_text_box(slide, 1.5, 2.5, 10, 1.2,
             "谢谢",
             font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 4.0, 10, 0.8,
             "Cross-Disciplinary Knowledge Extraction & Benchmarking System",
             font_size=18, color=RGBColor(0x88, 0xAA, 0xDD), align=PP_ALIGN.CENTER)

# ── 保存 ──
output_path = "/ssd/wangyuyang/git/benchmark/outputs/项目运行流程详解.pptx"
prs.save(output_path)
print(f"PPT 已保存至: {output_path}")
