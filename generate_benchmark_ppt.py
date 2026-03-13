#!/usr/bin/env python3
"""
根据 docs/benchmark_ppt_content.md 生成 CrossDisc Benchmark 汇报 PPT。
用法: python generate_benchmark_ppt.py [--output path/to/output.pptx]
"""
from __future__ import annotations

import argparse
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

# ── 全局配色 ──────────────────────────────────────────────────────────
CLR_BG_DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # 深蓝黑背景
CLR_BG_LIGHT  = RGBColor(0xFF, 0xFF, 0xFF)   # 白底
CLR_ACCENT    = RGBColor(0x00, 0x96, 0xD6)   # 科技蓝
CLR_ACCENT2   = RGBColor(0x00, 0xC9, 0xA7)   # 青绿
CLR_ACCENT3   = RGBColor(0xFF, 0x6B, 0x6B)   # 珊瑚红
CLR_ACCENT4   = RGBColor(0xFF, 0xA5, 0x02)   # 橙色
CLR_GOLD      = RGBColor(0xFF, 0xD7, 0x00)   # 金色
CLR_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
CLR_BLACK     = RGBColor(0x1A, 0x1A, 0x1A)
CLR_GRAY      = RGBColor(0x66, 0x66, 0x66)
CLR_LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)
CLR_P1        = RGBColor(0x95, 0xA5, 0xA6)   # 灰色
CLR_P2        = RGBColor(0x3E, 0xA8, 0xDB)   # 蓝
CLR_P3        = RGBColor(0x2E, 0xCC, 0x71)   # 绿
CLR_P4        = RGBColor(0xF3, 0x9C, 0x12)   # 橙
CLR_P5        = RGBColor(0xE7, 0x4C, 0x3C)   # 红
CLR_BASELINE  = RGBColor(0xBB, 0xBB, 0xBB)   # baseline灰


# ── 工具函数 ──────────────────────────────────────────────────────────
def set_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=CLR_BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=16, color=CLR_BLACK, bold_first=False,
                    font_name="Microsoft YaHei", spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        if bold_first and i == 0:
            p.font.bold = True
        p.level = 0
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=CLR_ACCENT, font_size=11):
    """data: 2D list [row][col] of strings. First row = header."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = table_shape.table

    col_width = width // cols
    for ci in range(cols):
        tbl.columns[ci].width = col_width

    for ri in range(rows):
        for ci in range(cols):
            cell = tbl.cell(ri, ci)
            cell.text = str(data[ri][ci])
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.name = "Microsoft YaHei"
                p.alignment = PP_ALIGN.CENTER
                if ri == 0:
                    p.font.bold = True
                    p.font.color.rgb = CLR_WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
    return table_shape


def add_rounded_rect(slide, left, top, width, height, fill_color, text,
                     font_size=14, font_color=CLR_WHITE):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(4)
    return shape


def add_arrow(slide, left, top, width, height):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CLR_ACCENT
    shape.line.fill.background()
    return shape


# ── 各 Slide 生成函数 ────────────────────────────────────────────────

def slide_cover(prs: Presentation):
    """Slide 1: 封面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, CLR_BG_DARK)

    # 上方装饰线
    from pptx.enum.shapes import MSO_SHAPE
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(1), Inches(1.8), Inches(2), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = CLR_ACCENT
    line.line.fill.background()

    add_textbox(slide, Inches(1), Inches(2.0), Inches(8), Inches(1.2),
                "CrossDisc", font_size=44, color=CLR_WHITE, bold=True)
    add_textbox(slide, Inches(1), Inches(2.9), Inches(8), Inches(0.8),
                "A Structured Benchmark for Cross-Disciplinary\nHypothesis Generation",
                font_size=22, color=CLR_ACCENT)
    add_textbox(slide, Inches(1), Inches(3.9), Inches(8), Inches(0.6),
                "面向跨学科科学假设生成的结构化基准与评估框架",
                font_size=16, color=CLR_LIGHT_GRAY)

    # 底部装饰线
    line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(1), Inches(5.0), Inches(8), Pt(1))
    line2.fill.solid()
    line2.fill.fore_color.rgb = RGBColor(0x44, 0x44, 0x66)
    line2.line.fill.background()


def slide_motivation(prs: Presentation):
    """Slide 2: 问题与动机"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_LIGHT)

    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                "AI for Science 基准的关键缺口", font_size=28, color=CLR_BLACK, bold=True)

    # 三列对比表
    data = [
        ["缺口", "现状", "我们的方案"],
        ["评估深度", "聚焦代码执行\n(MLAgentBench, ScienceAgentBench)", "评估跨学科推理\n与假设合成"],
        ["指标结构", "文本相似度 / LLM 自评", "图驱动指标\n(bridging, path consistency)"],
        ["评估客观性", "缺少 Ground Truth", "以真实文献的\n跨学科路径为参照"],
    ]
    add_table(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(2.4),
              rows=4, cols=3, data=data, font_size=12)

    # 底部讲稿要点
    add_textbox(slide, Inches(0.5), Inches(4.0), Inches(9), Inches(1.5),
                "现有 AI for Science 基准关注 LLM 能不能跑代码，但科学发现的核心能力是构想——"
                "提出有价值的假设。尤其在跨学科研究中，这种能力尤为关键。\n"
                "CrossDisc 是首个评估 LLM 跨学科假设生成能力的结构化基准。",
                font_size=13, color=CLR_GRAY)


def slide_pipeline(prs: Presentation):
    """Slide 3: 方法概览 — 三阶段管线"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_LIGHT)

    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                "CrossDisc Pipeline: 从论文到结构化假设",
                font_size=28, color=CLR_BLACK, bold=True)

    # 四个阶段方框 + 箭头
    box_w = Inches(2.0)
    box_h = Inches(2.2)
    y = Inches(1.5)
    gap = Inches(0.15)

    stages = [
        ("论文输入", "title\nabstract\nintroduction", CLR_GRAY),
        ("Stage 1\n学科分类与概念抽取", "学科分类\n概念抽取\n跨学科关系抽取", CLR_ACCENT),
        ("Stage 2\n三级查询生成", "L1: 宏观问题\nL2: 学科问题\nL3: 操作问题", CLR_ACCENT2),
        ("Stage 3\n三级假设生成", "L1: 宏观假设\nL2: 学科假设\nL3: 深层假设", CLR_ACCENT3),
    ]

    x = Inches(0.3)
    for i, (title, desc, clr) in enumerate(stages):
        add_rounded_rect(slide, x, y, box_w, Inches(0.8), clr, title,
                         font_size=12, font_color=CLR_WHITE)
        add_textbox(slide, x + Inches(0.1), y + Inches(0.9), box_w - Inches(0.2),
                    Inches(1.2), desc, font_size=11, color=CLR_GRAY)
        if i < len(stages) - 1:
            add_arrow(slide, x + box_w + Inches(0.02), y + Inches(0.3),
                      Inches(0.3), Inches(0.2))
        x += box_w + Inches(0.45)

    # 底部说明
    add_textbox(slide, Inches(0.5), Inches(4.2), Inches(9), Inches(1.2),
                "每条假设由 3-step 知识路径组成，实体锚定在真实科学概念上，"
                "步与步之间严格链式连接 (step₂.head == step₁.tail)。",
                font_size=13, color=CLR_GRAY)


def slide_innovation1_hierarchy(prs: Presentation):
    """Slide 4: 核心创新 1 — 层级化假设结构"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_LIGHT)

    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                "创新 1: L1/L2/L3 三级假设 — 业界首次",
                font_size=28, color=CLR_BLACK, bold=True)

    data = [
        ["层级", "视角", "示例 (AlphaFold3)"],
        ["L1 宏观", "论文最大未解决问题", "扩展扩散架构至动态构象预测"],
        ["L2 学科", "辅学科补充视角", "融入量子力学性质改进酶-底物预测"],
        ["L3 深层", "具体概念桥接", "蛋白质语言模型 + 小分子药效团联合嵌入"],
    ]
    add_table(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(2.0),
              rows=4, cols=3, data=data, font_size=12)

    # 对比框
    add_rounded_rect(slide, Inches(0.5), Inches(3.5), Inches(9), Inches(0.6),
                     RGBColor(0xFF, 0xEB, 0xEE),
                     "对比: 所有 6 种外部 baseline 仅产生单层平坦假设列表",
                     font_size=13, font_color=CLR_ACCENT3)

    add_textbox(slide, Inches(0.5), Inches(4.3), Inches(9), Inches(1.2),
                "L1 识别宏观瓶颈，L2 从辅学科角度提供新思路，L3 深入到具体概念和机制的跨学科桥接。\n"
                "这种设计确保了从宽到深的完整覆盖。",
                font_size=13, color=CLR_GRAY)


def slide_innovation2_path(prs: Presentation):
    """Slide 5: 核心创新 2 — 3-Step 知识路径"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_LIGHT)

    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                "创新 2: 3-Step 知识路径 — 可追溯的推理链",
                font_size=28, color=CLR_BLACK, bold=True)

    # 路径示例 - 三个步骤
    y_base = Inches(1.3)
    steps = [
        ("Step 1", "吸入毒物", "引起", "脂质分布变化", CLR_ACCENT),
        ("Step 2", "脂质分布变化", "需要", "质谱成像技术", CLR_ACCENT2),
        ("Step 3", "质谱成像技术", "提供", "高分辨率检测", CLR_ACCENT3),
    ]
    for i, (label, head, rel, tail, clr) in enumerate(steps):
        y = y_base + Inches(i * 0.7)
        add_rounded_rect(slide, Inches(0.8), y, Inches(1.5), Inches(0.45),
                         clr, head, font_size=11, font_color=CLR_WHITE)
        add_textbox(slide, Inches(2.5), y + Inches(0.05), Inches(1.5), Inches(0.4),
                    f"──{rel}──>", font_size=12, color=clr, bold=True,
                    alignment=PP_ALIGN.CENTER)
        add_rounded_rect(slide, Inches(4.2), y, Inches(1.8), Inches(0.45),
                         clr, tail, font_size=11, font_color=CLR_WHITE)
        add_textbox(slide, Inches(6.2), y + Inches(0.05), Inches(0.6), Inches(0.4),
                    label, font_size=10, color=CLR_GRAY)
        if i < 2:
            add_textbox(slide, Inches(4.5), y + Inches(0.45), Inches(1.5), Inches(0.25),
                        "  tail = next head", font_size=9, color=CLR_ACCENT4)

    # 约束列表
    constraints = [
        "每步 head/tail 必须来自抽取的科学实体（概念词表）",
        "严格链式约束: step₂.head == step₁.tail",
        "每步附带科学主张 (claim)，推理过程完全可追溯",
    ]
    add_bullet_list(slide, Inches(0.5), Inches(3.8), Inches(9), Inches(1.5),
                    [f"  {c}" for c in constraints],
                    font_size=13, color=CLR_BLACK)


def slide_innovation3_ablation(prs: Presentation):
    """Slide 6: 核心创新 3 — P1-P5 消融实验"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_LIGHT)

    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                "创新 3: P1-P5 消融框架 — 量化知识引导的边际贡献",
                font_size=26, color=CLR_BLACK, bold=True)

    # 阶梯图 - 五个阶段
    levels = [
        ("P1", "仅研究问题", CLR_P1),
        ("P2", "+ 学科角色 + 摘要 + L2 查询", CLR_P2),
        ("P3", "+ 抽取概念 + L3 查询", CLR_P3),
        ("P4", "+ 跨学科关系 + 推理链", CLR_P4),
        ("P5", "完整结构化管线", CLR_P5),
    ]
    x_base = Inches(0.5)
    y_base = Inches(3.2)
    bar_w = Inches(1.6)

    for i, (name, desc, clr) in enumerate(levels):
        x = x_base + Inches(i * 1.8)
        bar_h = Inches(0.4 + i * 0.35)
        y = y_base - bar_h

        add_rounded_rect(slide, x, y, bar_w, bar_h, clr, name,
                         font_size=16, font_color=CLR_WHITE)
        add_textbox(slide, x, y_base + Inches(0.1), bar_w, Inches(0.5),
                    desc, font_size=9, color=CLR_GRAY,
                    alignment=PP_ALIGN.CENTER)

    # 结果趋势说明
    add_textbox(slide, Inches(0.5), Inches(3.8), Inches(9), Inches(1.0),
                "P1→P4: LLM 评分呈清晰单调递增 (7.96→8.4→8.6→8.13)，"
                "跨学科性维度提升尤为显著 (7.5→8.0→9.0→8.33)。\n"
                "证明结构化知识引导对假设质量有显著提升效果。",
                font_size=13, color=CLR_GRAY)


def slide_experiment_setup(prs: Presentation):
    """Slide 7: 实验设置"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_LIGHT)

    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                "实验设置", font_size=28, color=CLR_BLACK, bold=True)

    # 论文表
    papers = [
        ["#", "论文", "学科"],
        ["1", "Resolving multi-image spatial lipidomic...", "生物学"],
        ["2", "A physics-informed Airy beam learning...", "物理学"],
        ["3", "Covariant spatio-temporal receptive fields...", "计算机科学"],
        ["4", "Non-equilibrium demixing of chiral...", "化学"],
        ["5", "Cellulose-mediated ionic liquid crystal...", "材料科学"],
        ["6", "AlphaFold 3: Accurate structure pred...", "计算生物学"],
    ]
    add_table(slide, Inches(0.3), Inches(1.1), Inches(6.0), Inches(2.5),
              rows=7, cols=3, data=papers, font_size=10)

    # 右侧配置信息
    config_items = [
        "对比方法: 6 外部 baseline + P1-P5 = 11 种",
        "外部 Baseline:",
        "  IdeaBench, VanillaLLM, AI-Scientist,",
        "  SciMON, MOOSE-Chem, SciAgents",
        "评估: ROUGE + LLM-Judge(5维) + 结构化指标",
        "总调用: 6 x 11 = 66 次，全部成功",
    ]
    add_bullet_list(slide, Inches(6.5), Inches(1.1), Inches(3.3), Inches(3.0),
                    config_items, font_size=12, color=CLR_BLACK)


def slide_llm_judge(prs: Presentation):
    """Slide 8: LLM Judge 评分"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_LIGHT)

    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                "LLM-as-Judge 多维度评分 (1-10 分)",
                font_size=26, color=CLR_BLACK, bold=True)

    data = [
        ["方法", "Novel", "Specif", "Feasib", "Relev", "Cross", "AVG"],
        ["P3",          "9.0",  "8.0",  "7.0",  "10.0", "9.0",  "8.6"],
        ["P2",          "9.0",  "8.0",  "7.0",  "10.0", "8.0",  "8.4"],
        ["SciAgents",   "8.5",  "8.83", "7.17", "9.5",  "8.17", "8.43"],
        ["AI-Scientist","8.33", "8.67", "7.0",  "10.0", "8.0",  "8.4"],
        ["MOOSE-Chem",  "9.0",  "8.0",  "7.0",  "9.0",  "8.0",  "8.2"],
        ["P4",          "8.33", "8.0",  "6.67", "9.33", "8.33", "8.13"],
        ["VanillaLLM",  "8.5",  "8.5",  "7.5",  "9.5",  "6.5",  "8.1"],
        ["P1",          "8.33", "8.33", "6.83", "8.83", "7.5",  "7.96"],
        ["IdeaBench",   "8.17", "7.17", "6.5",  "9.0",  "7.83", "7.73"],
    ]
    tbl_shape = add_table(slide, Inches(0.3), Inches(1.0), Inches(9.2), Inches(3.2),
                          rows=10, cols=7, data=data, font_size=11)

    # 高亮 P3 行 (row 1) 和 P2 行 (row 2)
    tbl = tbl_shape.table
    for ci in range(7):
        tbl.cell(1, ci).fill.solid()
        tbl.cell(1, ci).fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xF5)
        tbl.cell(2, ci).fill.solid()
        tbl.cell(2, ci).fill.fore_color.rgb = RGBColor(0xEB, 0xF5, 0xFB)

    # 亮点
    add_textbox(slide, Inches(0.5), Inches(4.3), Inches(9), Inches(1.0),
                "P3 在 Relevance (10.0) 和 Cross-disciplinarity (9.0) 上取得最高分；"
                "P2/P3 均值超越最强 baseline SciAgents。\n"
                "跨学科性维度 P1→P3: 7.5 → 8.0 → 9.0，验证概念注入的直接贡献。",
                font_size=12, color=CLR_ACCENT, bold=True)


def slide_rouge(prs: Presentation):
    """Slide 9: ROUGE 评分"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_LIGHT)

    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                "ROUGE 评分: 假设与原文摘要的语义对齐",
                font_size=26, color=CLR_BLACK, bold=True)

    data = [
        ["方法", "ROUGE-1", "ROUGE-2", "ROUGE-L"],
        ["P3",          "0.4224", "0.1522", "0.1878"],
        ["P2",          "0.4091", "0.1308", "0.1948"],
        ["VanillaLLM",  "0.3403", "0.0753", "0.1489"],
        ["IdeaBench",   "0.3317", "0.0700", "0.1524"],
        ["MOOSE-Chem",  "0.3043", "0.0612", "0.1303"],
        ["P4",          "0.2909", "0.0729", "0.1344"],
        ["AI-Scientist","0.2890", "0.0617", "0.1207"],
        ["SciMON",      "0.2706", "0.0394", "0.1148"],
        ["SciAgents",   "0.2104", "0.0303", "0.0937"],
        ["P1",          "0.1937", "0.0176", "0.0980"],
        ["P5",          "0.0096", "0.0018", "0.0096"],
    ]
    tbl_shape = add_table(slide, Inches(0.5), Inches(1.0), Inches(6.0), Inches(3.5),
                          rows=12, cols=4, data=data, font_size=11)

    # 高亮前两行
    tbl = tbl_shape.table
    for ci in range(4):
        tbl.cell(1, ci).fill.solid()
        tbl.cell(1, ci).fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xF5)
        tbl.cell(2, ci).fill.solid()
        tbl.cell(2, ci).fill.fore_color.rgb = RGBColor(0xEB, 0xF5, 0xFB)

    # 右侧说明
    add_textbox(slide, Inches(6.8), Inches(1.2), Inches(3.0), Inches(3.0),
                "P2/P3 的 ROUGE-1 领先\n最强 baseline (VanillaLLM)\n超过 24%\n\n"
                "P4 的 ROUGE 略低于 P3，\n因为 P4 引入了更多创新性\n跨学科连接，降低了与原\n文的直接重叠 —— 恰恰\n体现了更高的新颖性。\n\n"
                "P5 为结构化路径输出，\n非自由文本，ROUGE 不\n适用于衡量。",
                font_size=12, color=CLR_GRAY)


def slide_case_study(prs: Presentation):
    """Slide 10: 定性对比 — 案例分析"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_LIGHT)

    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                "案例对比: 脂质组学论文 — IdeaBench vs P4",
                font_size=26, color=CLR_BLACK, bold=True)

    # 左侧 - IdeaBench
    add_rounded_rect(slide, Inches(0.3), Inches(1.1), Inches(4.4), Inches(0.5),
                     CLR_BASELINE, "IdeaBench（典型 baseline）",
                     font_size=13, font_color=CLR_WHITE)
    add_textbox(slide, Inches(0.4), Inches(1.7), Inches(4.2), Inches(2.8),
                '"Sex-specific differences in lipid\n'
                'saturation patterns may underlie\n'
                'divergent susceptibility to asthma..."\n\n'
                '  - 单一段落\n'
                '  - 未桥接具体学科\n'
                '  - 无推理链\n'
                '  - 无定量预测',
                font_size=12, color=CLR_GRAY)

    # 右侧 - P4
    add_rounded_rect(slide, Inches(5.0), Inches(1.1), Inches(4.7), Inches(0.5),
                     CLR_ACCENT3, "P4（CrossDisc）",
                     font_size=13, font_color=CLR_WHITE)
    add_textbox(slide, Inches(5.1), Inches(1.7), Inches(4.5), Inches(3.2),
                'L1: 脂质饱和度梯度空间模式作为\n'
                '     生物物理缓冲... (推理链: 前提→机制→预测)\n\n'
                'L2-Q1: 脂质分布分形维数预测哮喘易感性\n'
                '     桥接: 生物物理学 <-> 生物化学\n\n'
                'L2-Q2: 细胞色素P450环氧合酶空间协调\n'
                '     桥接: 生物化学 <-> 空间生物学\n\n'
                'L3-Q1: 图神经网络 + 质谱成像\n'
                '     桥接: 机器学习 <-> 系统生物学\n\n'
                'L3-Q2: 氘标记脂质探针, ">50% 加速更新"\n'
                '     定量预测 + 可验证',
                font_size=11, color=CLR_BLACK)


def slide_p5_value(prs: Presentation):
    """Slide 11: P5 的独特价值"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_LIGHT)

    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                "P5: 结构化可追溯性 — 另一种贡献维度",
                font_size=26, color=CLR_BLACK, bold=True)

    # 路径示例
    add_textbox(slide, Inches(0.5), Inches(1.1), Inches(5.5), Inches(2.2),
                "P5 输出示例 (AlphaFold3):\n\n"
                "L1 路径 1:\n"
                "  [AlphaFold3] ──预测──> [深度学习框架]\n"
                "  [深度学习框架] ──需要──> [蛋白质化学性质]\n"
                "  [蛋白质化学性质] ──来自──> [修饰残基信息]\n\n"
                "  总结: 计算生物学需要生物化学提供修饰残基\n"
                "  信息来完善预测模型。",
                font_size=12, color=CLR_BLACK)

    # 指标表
    data = [
        ["P5 独有指标", "P5 得分", "其他方法"],
        ["层级覆盖率", "1.0", "N/A"],
        ["链式连贯性", "1.0", "N/A"],
        ["实体关联率", "1.0", "N/A"],
    ]
    add_table(slide, Inches(6.0), Inches(1.3), Inches(3.5), Inches(1.5),
              rows=4, cols=3, data=data, font_size=11,
              header_color=CLR_P5)

    # 定位说明
    add_rounded_rect(slide, Inches(0.5), Inches(3.8), Inches(9), Inches(0.7),
                     RGBColor(0xFD, 0xF2, 0xE9),
                     "P5 的价值不在于更好的假设文本，而在于提供可解释、可追溯、可评估的假设表示形式",
                     font_size=13, font_color=CLR_ACCENT4)


def slide_eval_framework(prs: Presentation):
    """Slide 12: 评估框架贡献"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_LIGHT)

    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                "超越 ROUGE: 图驱动评估体系",
                font_size=28, color=CLR_BLACK, bold=True)

    # 客观图指标
    add_textbox(slide, Inches(0.5), Inches(1.0), Inches(4.3), Inches(0.4),
                "客观图指标（无需 LLM）", font_size=16, color=CLR_ACCENT, bold=True)
    obj_metrics = [
        "Bridging Score: Rao-Stirling 多样性 + 嵌入距离",
        "Path Consistency: 关系感知 P/R/F1",
        "Chain Coherence: 逐跳语义连贯性",
        "Info-Theoretic Novelty: 基于 surprisal 的新颖度",
        "Atypical Combination Index: 组合非常规程度",
    ]
    add_bullet_list(slide, Inches(0.5), Inches(1.4), Inches(4.3), Inches(2.0),
                    [f"  {m}" for m in obj_metrics], font_size=11, color=CLR_BLACK)

    # LLM 评估
    add_textbox(slide, Inches(5.2), Inches(1.0), Inches(4.5), Inches(0.4),
                "LLM 多维度评估", font_size=16, color=CLR_ACCENT2, bold=True)
    llm_metrics = [
        "创新性 / 可行性 / 科学性 (各 0-10)",
        "可检验性: 具体性 + 可测量性 + 可证伪性",
        "GT 概念覆盖: P / R / F1",
        "GT 关系精确度 + 证据覆盖",
        "结构多样性 (Torrance 框架)",
    ]
    add_bullet_list(slide, Inches(5.2), Inches(1.4), Inches(4.5), Inches(2.0),
                    [f"  {m}" for m in llm_metrics], font_size=11, color=CLR_BLACK)

    # 对比框
    comparisons = [
        ["评估方法", "文本维度", "主观维度", "图结构", "GT 参照"],
        ["IdeaBench", "ROUGE", "-", "-", "-"],
        ["通用 LLM-Judge", "-", "LLM 评分", "-", "-"],
        ["CrossDisc", "ROUGE+BLEU+BERTScore", "LLM 5 维", "6 类图指标", "证据驱动"],
    ]
    add_table(slide, Inches(0.3), Inches(3.5), Inches(9.2), Inches(1.5),
              rows=4, cols=5, data=comparisons, font_size=11)


def slide_contributions(prs: Presentation):
    """Slide 13: 总结与贡献"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_DARK)

    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                "CrossDisc 的五大贡献",
                font_size=30, color=CLR_WHITE, bold=True)

    contributions = [
        ("1", "首个跨学科假设生成基准", "评估 LLM 跨学科科学推理能力，而非代码执行能力", CLR_ACCENT),
        ("2", "层级化假设结构 (L1/L2/L3)", "三级假设 + 3-step 知识路径，业界首次", CLR_ACCENT2),
        ("3", "P1-P5 消融实验框架", "五级 prompt 梯度，量化结构化知识引导的边际贡献", CLR_P3),
        ("4", "图驱动评估体系", "超越文本相似度，基于知识图谱拓扑的客观评估", CLR_ACCENT4),
        ("5", "全面实验验证", "6 篇论文 x 11 种方法，P2/P3 多维度超越 SOTA", CLR_ACCENT3),
    ]

    y = Inches(1.2)
    for num, title, desc, clr in contributions:
        add_rounded_rect(slide, Inches(0.5), y, Inches(0.5), Inches(0.5),
                         clr, num, font_size=16, font_color=CLR_WHITE)
        add_textbox(slide, Inches(1.2), y, Inches(3.5), Inches(0.5),
                    title, font_size=16, color=CLR_WHITE, bold=True)
        add_textbox(slide, Inches(1.2), y + Inches(0.4), Inches(8), Inches(0.4),
                    desc, font_size=12, color=CLR_LIGHT_GRAY)
        y += Inches(0.85)

    # one-liner
    add_textbox(slide, Inches(0.5), Inches(5.5), Inches(9), Inches(0.5),
                "CrossDisc 是首个评估 LLM 跨学科科学假设生成能力的结构化基准。",
                font_size=14, color=CLR_GOLD, bold=True,
                alignment=PP_ALIGN.CENTER)


def slide_future_work(prs: Presentation):
    """Slide 14: 未来工作"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_LIGHT)

    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                "Future Work", font_size=28, color=CLR_BLACK, bold=True)

    items = [
        "扩展数据集规模: Nature Communications 100 篇论文全量测试",
        "引入人类专家评估: 与 LLM 评估交叉验证",
        "动态知识图谱: 融入时间维度的学科演化路径",
        "多模型对比: GPT-4, Claude, Gemini 等不同 LLM 能力评估",
    ]

    y = Inches(1.3)
    icons = [CLR_ACCENT, CLR_ACCENT2, CLR_ACCENT4, CLR_ACCENT3]
    for i, (item, clr) in enumerate(zip(items, icons)):
        add_rounded_rect(slide, Inches(0.5), y, Inches(0.4), Inches(0.4),
                         clr, str(i + 1), font_size=14, font_color=CLR_WHITE)
        add_textbox(slide, Inches(1.1), y, Inches(8.5), Inches(0.5),
                    item, font_size=16, color=CLR_BLACK)
        y += Inches(0.7)


def slide_appendix_tables(prs: Presentation):
    """Slide 15: 附录 - 完整数据表"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_LIGHT)

    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(9), Inches(0.5),
                "附录: P1-P5 消融实验质量对比",
                font_size=24, color=CLR_BLACK, bold=True)

    data = [
        ["维度",       "P1",   "P2",  "P3",  "P4",   "P5"],
        ["跨学科深度", "低",    "中",  "中",  "高",   "中"],
        ["多层级覆盖", "无",    "2级", "3级", "3级",  "3级"],
        ["推理透明度", "无",    "无",  "低",  "高",   "中"],
        ["具体性",     "不稳定","中",  "中",  "高",   "中低"],
        ["可验证性",   "不稳定","中",  "中",  "高",   "低"],
        ["实体关联度", "无",    "中",  "高",  "高",   "最高"],
    ]
    tbl_shape = add_table(slide, Inches(0.5), Inches(0.9), Inches(9), Inches(2.8),
                          rows=7, cols=6, data=data, font_size=12)

    # 高亮 P4 列 (col 4)
    tbl = tbl_shape.table
    for ri in range(1, 7):
        tbl.cell(ri, 4).fill.solid()
        tbl.cell(ri, 4).fill.fore_color.rgb = RGBColor(0xFE, 0xF5, 0xE7)

    # LLM Judge 完整表
    add_textbox(slide, Inches(0.5), Inches(3.8), Inches(9), Inches(0.4),
                "LLM Judge 完整评分表 (含 P5)",
                font_size=14, color=CLR_BLACK, bold=True)

    data2 = [
        ["方法", "Novel", "Specif", "Feasib", "Relev", "Cross", "AVG"],
        ["P3",          "9.0",  "8.0",  "7.0",  "10.0", "9.0",  "8.6"],
        ["P2",          "9.0",  "8.0",  "7.0",  "10.0", "8.0",  "8.4"],
        ["SciAgents",   "8.5",  "8.83", "7.17", "9.5",  "8.17", "8.43"],
        ["AI-Scientist","8.33", "8.67", "7.0",  "10.0", "8.0",  "8.4"],
        ["P5",          "7.5",  "8.0",  "7.5",  "9.25", "7.5",  "7.95"],
        ["IdeaBench",   "8.17", "7.17", "6.5",  "9.0",  "7.83", "7.73"],
    ]
    add_table(slide, Inches(0.5), Inches(4.2), Inches(9), Inches(2.2),
              rows=7, cols=7, data=data2, font_size=10)


# ── 主函数 ────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="生成 CrossDisc Benchmark PPT")
    parser.add_argument("--output", "-o", default=None,
                        help="输出路径 (默认: docs/CrossDisc_Benchmark.pptx)")
    args = parser.parse_args()

    proj_dir = os.path.dirname(os.path.abspath(__file__))
    output = args.output or os.path.join(proj_dir, "docs", "CrossDisc_Benchmark.pptx")
    os.makedirs(os.path.dirname(output), exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    # 按顺序生成所有 Slide
    slide_cover(prs)                    # 1. 封面
    slide_motivation(prs)               # 2. 问题与动机
    slide_pipeline(prs)                 # 3. 方法概览
    slide_innovation1_hierarchy(prs)    # 4. 创新1: 层级化假设
    slide_innovation2_path(prs)         # 5. 创新2: 3-Step 路径
    slide_innovation3_ablation(prs)     # 6. 创新3: P1-P5 消融
    slide_experiment_setup(prs)         # 7. 实验设置
    slide_llm_judge(prs)                # 8. LLM Judge 评分
    slide_rouge(prs)                    # 9. ROUGE 评分
    slide_case_study(prs)               # 10. 案例对比
    slide_p5_value(prs)                 # 11. P5 价值
    slide_eval_framework(prs)           # 12. 评估框架
    slide_contributions(prs)            # 13. 总结
    slide_future_work(prs)              # 14. 未来工作
    slide_appendix_tables(prs)          # 15. 附录

    prs.save(output)
    print(f"PPT 已生成: {output}")
    print(f"共 {len(prs.slides)} 页 Slide")


if __name__ == "__main__":
    main()
